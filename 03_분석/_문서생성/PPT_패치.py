# -*- coding: utf-8 -*-
"""
PPT 패치 도구 — 사용자가 편집 중인 제출본을 '원본'으로 두고 그 위에 수정만 얹는다.

배경: 04_제출물의 PPT는 사용자가 PowerPoint로 직접 편집한다. 생성기를 다시 돌려
      통째로 덮어쓰면 그 편집이 전부 사라진다. 그래서 수정은 항상
      "현재 제출본을 읽어 → 필요한 부분만 고쳐 → 새 파일로 저장" 순서로 한다.

원칙 3가지
  1. 제출본에 직접 쓰지 않는다. 결과는 항상 staging(05_작업자료/_PPT생성본)에 저장한다.
  2. 반영(deploy)은 사용자가 명시적으로 요청할 때만, 그리고 파일이 열려 있지 않을 때만 한다.
  3. PowerPoint COM은 쓰지 않는다. 실행 중인 인스턴스에 붙어 사용자의 창을 닫아버린다.

사용 예:
    from PPT_패치 import load, replace_text, save, is_open

    prs = load('결과')                      # 현재 제출본을 읽어온다 (원본은 건드리지 않음)
    n = replace_text(prs, '옛 문구', '새 문구')
    save(prs, '결과')                        # staging에 저장, 경로를 돌려준다
"""
import os
import shutil

from pptx import Presentation

__all__ = ['load', 'save', 'deploy', 'is_open', 'replace_text', 'find_text',
           'set_paragraph', 'dump', 'SUBMIT', 'STAGING']


def _project_root(start):
    d = os.path.dirname(os.path.abspath(start))
    while d != os.path.dirname(d):
        if os.path.isdir(os.path.join(d, '04_제출물')):
            return d
        d = os.path.dirname(d)
    raise RuntimeError('프로젝트 루트를 찾지 못했습니다')


BASE = _project_root(__file__)
SUBMIT = os.path.join(BASE, '04_제출물')
STAGING = os.path.join(BASE, '05_작업자료', '_PPT생성본')
DECKS = {'결과': '데이터분석_결과보고서.pptx', '계획서': '데이터분석_계획서.pptx'}

# 다른 사람이 편집 중이라 손대면 안 되는 덱. 대표자가 넘겨줄 때 이 줄을 비운다.
# 읽기(load)는 허용하고 쓰기(deploy)만 막는다 — 확인은 해야 하니까.
HOLD = {'결과': '팀원이 작업 중 (2026-08-12~). 대표자가 다시 넘겨줄 때까지 반영 금지'}


def _name(deck):
    if deck in DECKS:
        return DECKS[deck]
    if deck.endswith('.pptx'):
        return deck
    raise KeyError(f'알 수 없는 덱: {deck} (가능: {list(DECKS)})')


def is_open(deck):
    """그 파일에 지금 쓸 수 없는 상태인지 판정한다.

    ~$ 잠금 파일만 보면 안 된다 — PowerPoint가 비정상 종료하면 잠금 파일이 남아
    닫혀 있는데도 열린 것으로 오판한다(실제로 이 때문에 계획서 작업이 여러 번 건너뛰어졌다).
    그래서 실제로 쓰기 모드로 열어 보는 것을 1차 판정으로 삼고, ~$ 는 참고만 한다.
    """
    path = os.path.join(SUBMIT, _name(deck))
    if not os.path.exists(path):
        return False
    try:
        with open(path, 'r+b'):        # 다른 프로세스가 잡고 있으면 PermissionError
            return False
    except PermissionError:
        return True
    except OSError:
        return True


def stale_lock(deck):
    """쓸 수 있는데도 ~$ 잠금 파일이 남아 있는 상태 — 비정상 종료의 잔여물."""
    n = _name(deck)
    return os.path.exists(os.path.join(SUBMIT, '~$' + n)) and not is_open(deck)


def load(deck):
    """현재 제출본을 읽어 Presentation 으로 돌려준다. 원본 파일은 열어만 보고 쓰지 않는다.

    편집 중이어도 안전하다 — 읽기는 잠금과 무관하다. 다만 사용자가 아직 저장하지 않은
    변경은 파일에 없으므로, 반영 전에 '저장했는지' 확인하는 것은 호출자의 몫이다.
    """
    src = os.path.join(SUBMIT, _name(deck))
    if not os.path.exists(src):
        raise FileNotFoundError(src)
    return Presentation(src)


def save(prs, deck, suffix=''):
    """staging 폴더에 저장하고 경로를 돌려준다. 제출본은 절대 건드리지 않는다."""
    os.makedirs(STAGING, exist_ok=True)
    n = _name(deck)
    if suffix:
        stem, ext = os.path.splitext(n)
        n = f'{stem}{suffix}{ext}'
    out = os.path.join(STAGING, n)
    prs.save(out)
    return out


def deploy(deck, suffix=''):
    """staging 결과를 제출본 자리로 옮긴다. 사용자가 명시적으로 요청했을 때만 호출할 것.

    파일이 열려 있으면 아무것도 하지 않고 이유를 돌려준다 — 열린 채로 덮어쓰면
    사용자가 PowerPoint에서 저장하는 순간 되돌아가거나, 편집분이 사라진다.
    """
    n = _name(deck)
    if deck in HOLD:
        return False, f'거부: {n} 은(는) 반영 보류 중입니다 — {HOLD[deck]}'
    if is_open(deck):
        return False, f'{n} 이(가) PowerPoint에서 열려 있습니다. 닫은 뒤 다시 시도하세요.'
    stem, ext = os.path.splitext(n)
    src = os.path.join(STAGING, f'{stem}{suffix}{ext}' if suffix else n)
    if not os.path.exists(src):
        return False, f'staging에 {os.path.basename(src)} 이(가) 없습니다.'
    dst = os.path.join(SUBMIT, n)
    backup = os.path.join(STAGING, f'{stem}_반영직전_백업{ext}')
    if os.path.exists(dst):
        shutil.copy2(dst, backup)      # 되돌릴 수 있도록 직전 상태를 남긴다
    shutil.copy2(src, dst)
    return True, f'{n} 반영 완료 (직전 상태는 {os.path.basename(backup)} 에 보관)'


# ── 내용 수정 ───────────────────────────────────────────────
def _frames(prs):
    """덱 안의 모든 텍스트 프레임을 (슬라이드번호, 프레임)으로 훑는다 (표 셀 포함)."""
    for i, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            if sh.has_text_frame:
                yield i, sh.text_frame
            if sh.has_table:
                for row in sh.table.rows:
                    for cell in row.cells:
                        yield i, cell.text_frame


def find_text(prs, needle):
    """needle 을 포함한 위치를 [(슬라이드번호, 문단 전체 텍스트)] 로 돌려준다."""
    hits = []
    for i, tf in _frames(prs):
        for p in tf.paragraphs:
            t = ''.join(r.text for r in p.runs)
            if needle in t:
                hits.append((i, t))
    return hits


def replace_text(prs, old, new, slide=None):
    """문단 단위로 old→new 치환. 서식은 그 문단 첫 run 의 것을 유지한다.

    한 문단이 여러 run 으로 쪼개져 있어도(한글 입력 중 흔하다) 문단을 합쳐서 비교하므로
    걸쳐 있는 문자열도 잡힌다. 바꾼 개수를 돌려준다.
    """
    n = 0
    for i, tf in _frames(prs):
        if slide is not None and i != slide:
            continue
        for p in tf.paragraphs:
            runs = p.runs
            if not runs:
                continue
            whole = ''.join(r.text for r in runs)
            if old not in whole:
                continue
            runs[0].text = whole.replace(old, new)
            for r in runs[1:]:
                r.text = ''
            n += 1
    return n


def set_paragraph(prs, slide, shape_idx, para_idx, text):
    """슬라이드의 특정 도형·문단 텍스트를 통째로 바꾼다 (서식 유지). 위치를 아는 경우용."""
    sh = prs.slides[slide - 1].shapes[shape_idx]
    p = sh.text_frame.paragraphs[para_idx]
    if not p.runs:
        raise ValueError('빈 문단이라 서식을 유지할 수 없습니다')
    p.runs[0].text = text
    for r in p.runs[1:]:
        r.text = ''


def insert_bullet_before(prs, anchor_text, new_text):
    """anchor_text 를 담은 문단 바로 앞에 같은 서식의 문단을 하나 끼워 넣는다.

    문단 XML을 통째로 복제하므로 글머리표·들여쓰기·글꼴이 그대로 유지된다.
    문장을 기존 문단에 이어 붙이면 한 줄로 뭉치므로, 항목을 늘릴 때는 이 함수를 쓴다.
    """
    import copy
    n = 0
    for _, tf in _frames(prs):
        for p in tf.paragraphs:
            if anchor_text not in ''.join(r.text for r in p.runs):
                continue
            new_p = copy.deepcopy(p._p)
            p._p.addprevious(new_p)
            from pptx.text.text import _Paragraph
            np = _Paragraph(new_p, p._parent)
            np.runs[0].text = new_text
            for r in np.runs[1:]:
                r.text = ''
            n += 1
    return n


def add_slide_at(prs, index):
    """빈 슬라이드를 만들어 index(0부터) 자리로 옮기고 돌려준다.

    python-pptx에는 순서 변경 API가 없어 sldIdLst를 직접 다룬다.
    기존 슬라이드는 그대로 두므로 사용자의 편집이 보존된다.
    """
    layout = prs.slide_layouts[6]          # 빈 레이아웃
    s = prs.slides.add_slide(layout)
    lst = prs.slides._sldIdLst
    el = lst[-1]
    lst.remove(el)
    lst.insert(index, el)
    return s


def renumber_pages(prs, start_slide=2, first_no=3):
    """슬라이드 우상단의 두 자리 쪽번호를 다시 매긴다 (슬라이드를 끼워 넣은 뒤 사용).

    헤더가 그리는 번호 텍스트박스는 가로 11.85in 근처에 있고 내용이 숫자뿐이라 그것만 고른다.
    start_slide는 0부터 센 인덱스, first_no는 그 슬라이드에 찍을 번호.
    """
    from pptx.util import Emu
    n = first_no
    changed = 0
    for s in list(prs.slides)[start_slide:]:
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text.strip()
            if t.isdigit() and len(t) <= 2 and Emu(sh.left).inches > 11.0:
                p = sh.text_frame.paragraphs[0]
                if p.runs:
                    p.runs[0].text = f'{n:02d}'
                    for r in p.runs[1:]:
                        r.text = ''
                    changed += 1
                break
        n += 1
    return changed


def dump(deck_or_prs):
    """슬라이드별 제목과 도형 수를 찍어 본다 — 패치 대상 찾을 때 쓴다."""
    prs = load(deck_or_prs) if isinstance(deck_or_prs, str) else deck_or_prs
    for i, s in enumerate(prs.slides, 1):
        head = ''
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                head = sh.text_frame.text.strip().splitlines()[0]
                break
        print(f'[{i:2d}] 도형 {len(s.shapes):3d}  {head[:60]}')


if __name__ == '__main__':
    import sys
    if (sys.stdout.encoding or '').lower() not in ('utf-8', 'utf8'):
        import io
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
    for k in DECKS:
        print(f'── {k} ({DECKS[k]}) · 열림={is_open(k)}')
        dump(k)
        print()
