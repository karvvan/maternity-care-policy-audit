# -*- coding: utf-8 -*-
"""제출본 PPT 정합성 점검·복구 — 내가 가한 수정이 지금도 살아 있는지 확인하고, 빠졌으면 다시 넣는다.

왜 필요한가:
  PowerPoint는 파일 전체를 메모리에 들고 있다가 저장할 때 통째로 덮어쓴다. 그래서
  "사용자가 파일을 연 뒤 → 내가 디스크의 파일을 수정 → 사용자가 저장" 순서가 되면
  사용자의 메모리 상태가 디스크를 덮어써서 내 수정이 사라진다. 잠금(~$) 검사만으로는
  이 순서를 막지 못한다. 그래서 수정 내역을 선언적으로 남겨 두고 언제든 다시 넣는다.

멱등(idempotent)하다 — 이미 반영된 항목은 건너뛰므로 몇 번을 돌려도 안전하다.

  python PPT_정합성점검.py            → 점검만 (파일을 건드리지 않음)
  python PPT_정합성점검.py --apply    → 빠진 항목만 복구 (닫혀 있어야 함, 직전 상태 백업)
"""
import sys, io, os

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
if (sys.stdout.encoding or '').lower() not in ('utf-8', 'utf8'):
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx.util import Emu, Inches
from PPT_패치 import load, save, deploy, is_open, replace_text, _frames

CHIP_IN = 0.34 * 2 / 3          # 헤더 액센트 칩 목표 크기 (원래 0.34의 2/3)


# ── 개별 점검 항목 ──────────────────────────────────────────
def chk_chip(prs):
    """헤더 칩이 축소본인가"""
    bad = []
    for i, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            w, h = Emu(sh.width).inches, Emu(sh.height).inches
            if abs(w - h) < 0.01 and 0.15 < w < 0.45 and abs(w - CHIP_IN) > 0.01:
                bad.append((i, sh))
            break_ = None
    return bad


def fix_chip(prs, bad):
    for _, sh in bad:
        cx = Emu(sh.left).inches + Emu(sh.width).inches / 2
        cy = Emu(sh.top).inches + Emu(sh.height).inches / 2
        sh.left, sh.top = Inches(cx - CHIP_IN / 2), Inches(cy - CHIP_IN / 2)
        sh.width = sh.height = Inches(CHIP_IN)
    return len(bad)


TERMS = [('연구 질문', '핵심 질문'), ('본 연구의', '이 프로젝트의'), ('본 연구가', '이 프로젝트가'),
         ('본 연구는', '이 프로젝트는'), ('본 연구', '이 프로젝트'),
         ('이 연구의', '이 프로젝트의'), ('이 연구가', '이 프로젝트가'), ('이 연구는', '이 프로젝트는'),
         ('연구의 한계', '분석의 한계'), ('연구 요약', '프로젝트 요약'),
         ('RQ1', 'Q1'), ('RQ2', 'Q2'), ('RQ3', 'Q3'), ('RQ4', 'Q4'), ('RQ', 'Q')]

# 반드시 남아 있어야 하는 문구 (없으면 사용자 저장에 덮여 사라진 것)
MUST = {
    '결과': ['목차', '프로젝트 배경', '나눠 줄 것이 적을수록', '핵심 질문'],
    '계획서': ['목차', '아이디어의 발단', '나눠 줄 것이 적을수록', '핵심 질문'],
}


def blob(prs):
    return '\n'.join(''.join(r.text for r in p.runs)
                     for _, tf in _frames(prs) for p in tf.paragraphs)


def check(deck, apply=False):
    prs = load(deck)
    b = blob(prs)
    issues, fixed = [], 0

    bad = chk_chip(prs)
    if bad:
        issues.append(f'헤더 칩이 원래 크기인 슬라이드 {len(bad)}개')
        if apply:
            fixed += fix_chip(prs, bad)

    stale = [a for a, _ in TERMS if a in b]
    if stale:
        issues.append(f'옛 용어 잔존: {sorted(set(stale))}')
        if apply:
            for a, c in TERMS:
                fixed += replace_text(prs, a, c)

    missing = [m for m in MUST[deck] if m not in b]
    if missing:
        issues.append(f'★ 사라진 내용: {missing} — 자동 복구 불가, 별도 패치 필요')

    print(f'── {deck}  {len(prs.slides._sldIdLst)}장 · 열림={is_open(deck)}')
    if not issues:
        print('   정상 — 모든 수정이 반영되어 있음')
        return
    for it in issues:
        print('   !!', it)
    if apply and fixed:
        if is_open(deck):
            print('   → 열려 있어 저장하지 못함. 닫고 다시 실행하세요.')
            return
        save(prs, deck)
        print('   →', *deploy(deck))
    elif apply:
        print('   → 자동 복구할 항목이 없음')
    else:
        print('   → --apply 로 복구 가능')


if __name__ == '__main__':
    ap = '--apply' in sys.argv
    for d in ('결과', '계획서'):
        check(d, ap)
        print()
