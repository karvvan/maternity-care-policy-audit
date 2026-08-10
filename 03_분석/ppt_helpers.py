# -*- coding: utf-8 -*-
"""
PPT 공용 헬퍼 v2 — 다크 '오로라' 테마
글래스모피즘 카드(반투명+그림자), 그라데이션, 3D 조형물(design_assets) 합성.
기존 덱 코드와의 호환을 위해 색 상수 이름(NAVY/DARK/RED/GRAY/LIGHT/...)은 유지하되
다크 테마에 맞는 역할로 재해석한다:
  DARK  = 본문 텍스트(밝은 회백)   NAVY = 카드 제목(앰버)   GRAY = 보조 텍스트
  LIGHT = 글래스 카드 채움(마커)   REDBG = 경고 글래스(마커)  RED = 경고 텍스트
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

import design_assets

ASSETS = design_assets.build_defaults()

# ── 팔레트 ──────────────────────────────────────────────────
DARK = RGBColor(0xEC, 0xF1, 0xF8)    # 본문 텍스트 (밝음)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA7, 0xB6, 0xC9)    # 보조 텍스트
NAVY = RGBColor(0xFF, 0xC9, 0x8A)    # 카드 제목 (앰버)
RED = RGBColor(0xFF, 0x8A, 0x7A)     # 경고 텍스트 (코랄)
TEAL = RGBColor(0x6E, 0xE7, 0xD2)
BLUE = RGBColor(0x8A, 0xB4, 0xFF)
LIGHT = RGBColor(0xF2, 0xF6, 0xFA)   # 마커: 글래스 카드
REDBG = RGBColor(0xFB, 0xF0, 0xEE)   # 마커: 경고 글래스 카드
LINE = RGBColor(0xB9, 0xC4, 0xCF)    # 마커: (무시됨)
INK = RGBColor(0x10, 0x1B, 0x2E)     # 어두운 잉크(밝은 배경 위 텍스트용)
FONT = '맑은 고딕'

_ACCENT1, _ACCENT2 = 'FF7A59', 'FFB86B'     # 코랄→앰버
_BLUE1, _BLUE2 = '2E4E86', '1B2C4F'         # 스탯 블록 그라데이션


def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs, prs.slide_layouts[6]


# ── 저수준 스타일 유틸 ──────────────────────────────────────
def _spPr(sp):
    return sp._element.spPr


def alpha_fill(sp, hex_str, alpha):
    """반투명 단색 채움 (alpha: 0~100)"""
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor.from_string(hex_str)
    clr = _spPr(sp).find(qn('a:solidFill')).find(qn('a:srgbClr'))
    a = etree.SubElement(clr, qn('a:alpha'))
    a.set('val', str(int(alpha * 1000)))


def grad_fill(sp, c1, c2, angle=90):
    sp.fill.gradient()
    stops = sp.fill.gradient_stops
    stops[0].color.rgb = RGBColor.from_string(c1)
    stops[0].position = 0.0
    stops[1].color.rgb = RGBColor.from_string(c2)
    stops[1].position = 1.0
    try:
        sp.fill.gradient_angle = angle
    except Exception:
        pass


def line_style(sp, hex_str, alpha=100, width=1.0):
    sp.line.color.rgb = RGBColor.from_string(hex_str)
    sp.line.width = Pt(width)
    if alpha < 100:
        ln = _spPr(sp).find(qn('a:ln'))
        clr = ln.find(qn('a:solidFill')).find(qn('a:srgbClr'))
        a = etree.SubElement(clr, qn('a:alpha'))
        a.set('val', str(int(alpha * 1000)))


def shadow(sp, blur=0.16, dist=0.055, alpha=52, direction=5400000):
    spPr = _spPr(sp)
    old = spPr.find(qn('a:effectLst'))
    if old is not None:
        spPr.remove(old)
    el = etree.SubElement(spPr, qn('a:effectLst'))
    sh = etree.SubElement(el, qn('a:outerShdw'))
    sh.set('blurRad', str(int(Inches(blur))))
    sh.set('dist', str(int(Inches(dist))))
    sh.set('dir', str(direction))
    sh.set('rotWithShape', '0')
    clr = etree.SubElement(sh, qn('a:srgbClr'))
    clr.set('val', '000000')
    a = etree.SubElement(clr, qn('a:alpha'))
    a.set('val', str(int(alpha * 1000)))


def cell_fill(cell, hex_str, alpha=100):
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ('a:solidFill', 'a:noFill', 'a:gradFill', 'a:pattFill'):
        for e in tcPr.findall(qn(tag)):
            tcPr.remove(e)
    fill = etree.SubElement(tcPr, qn('a:solidFill'))
    clr = etree.SubElement(fill, qn('a:srgbClr'))
    clr.set('val', hex_str)
    if alpha < 100:
        a = etree.SubElement(clr, qn('a:alpha'))
        a.set('val', str(int(alpha * 1000)))


# ── 배경·조형물 ─────────────────────────────────────────────
def bg(slide, kind='main'):
    """전면 배경 이미지 (main/title/closing) — 가장 먼저 호출할 것"""
    slide.shapes.add_picture(ASSETS[f'bg_{kind}'], 0, 0, Inches(13.333), Inches(7.5))


def glow(slide, x, y, w, h=0.055):
    """코랄-앰버 글로우 라인"""
    slide.shapes.add_picture(ASSETS['glow_bar'], Inches(x), Inches(y), Inches(w), Inches(h))


def deco(slide, key, x, y, w):
    """3D 조형물 배치 (sphere_coral/sphere_blue/sphere_teal/ring_blue/ring_coral)"""
    from PIL import Image
    iw, ih = Image.open(ASSETS[key]).size
    slide.shapes.add_picture(ASSETS[key], Inches(x), Inches(y), Inches(w), Inches(w * ih / iw))


# ── 텍스트 ──────────────────────────────────────────────────
def set_font(run, size, bold=False, color=DARK):
    f = run.font
    f.size, f.bold, f.name = Pt(size), bold, FONT
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def text_block(slide, x, y, w, h, items, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (t, size, bold, color, after) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(after)
        p.line_spacing = 1.14
        r = p.add_run(); r.text = t
        set_font(r, size, bold, color)
        if t.startswith('• '):
            pPr = p._p.get_or_add_pPr()
            pPr.set('marL', '160020'); pPr.set('indent', '-160020')
    return tb


# ── 도형 ────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill, line_color=None):
    """테마 인지 사각형: LIGHT→글래스, REDBG→경고 글래스, 그 외→그라데이션 스탯 블록"""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.adjustments[0] = min(0.5, 0.09 / min(w, h) if min(w, h) > 0.4 else 0.18)
    sp.shadow.inherit = False
    if fill == LIGHT:
        alpha_fill(sp, 'FFFFFF', 7)
        line_style(sp, 'FFFFFF', 22, 1.0)
        shadow(sp)
    elif fill == REDBG:
        alpha_fill(sp, 'FF7A59', 13)
        line_style(sp, 'FF9A7B', 32, 1.0)
        shadow(sp)
    elif fill == WHITE:
        sp.fill.solid(); sp.fill.fore_color.rgb = WHITE
        sp.line.fill.background()
        shadow(sp, blur=0.2, alpha=58)
    else:  # 스탯 블록·헤더 칩 — 남색 그라데이션
        grad_fill(sp, _BLUE1, _BLUE2, angle=115)
        line_style(sp, 'FFFFFF', 14, 0.75)
        shadow(sp)
    return sp


def accent_chip(slide, x, y, w, h):
    """코랄-앰버 그라데이션 칩"""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.adjustments[0] = 0.5
    sp.shadow.inherit = False
    grad_fill(sp, _ACCENT1, _ACCENT2, angle=45)
    sp.line.fill.background()
    shadow(sp, blur=0.14, dist=0.04, alpha=45)
    return sp


def card(slide, x, y, w, h, title, body, fill=LIGHT, tcolor=NAVY, tsize=13, bsize=11.5):
    rect(slide, x, y, w, h, fill)
    items = [(title, tsize, True, tcolor, 3.5)]
    for line in body:
        items.append((line if line.startswith(('→', '—')) else '• ' + line, bsize, False, DARK, 2))
    text_block(slide, x + 0.2, y + 0.14, w - 0.4, h - 0.28, items)


def header(slide, no, title, footer_label, sub=None):
    """다크 헤더: 액센트 칩 + 흰 제목 + 글로우 바 + 큰 페이지 번호"""
    bg(slide, 'main')
    accent_chip(slide, 0.55, 0.52, 0.34, 0.34)
    text_block(slide, 1.05, 0.3, 10.6, 0.75, [(title, 21.5, True, WHITE, 0)], anchor=MSO_ANCHOR.MIDDLE)
    slide.shapes.add_picture(ASSETS['glow_bar'], Inches(0.57), Inches(1.06), Inches(4.6), Inches(0.05))
    text_block(slide, 11.85, 0.18, 1.3, 0.85, [(f'{no:02d}', 30, True, RGBColor(0x3A, 0x4E, 0x71), 0)])
    if sub:
        text_block(slide, 0.6, 1.2, 12.2, 0.42, [(sub, 12.5, False, GRAY, 0)])
    text_block(slide, 0.55, 7.1, 12.3, 0.35,
               [(f'2026 명지대학교 창의적 SW프로그램 경진대회 · 빅데이터 분석 — {footer_label}', 8.5, False,
                 RGBColor(0x5E, 0x70, 0x8C), 0)])


def table(slide, x, y, w, rows, widths=None, header_size=11.5, body_size=11, row_h=0.42,
          highlight_rows=(), highlight_color=None):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(row_h * len(rows)))
    tbl = shape.table
    # 기본 표 스타일 제거(밴딩 끔)
    tbl.first_row = False
    tbl.horz_banding = False
    if widths:
        for j, wd in enumerate(widths):
            tbl.columns[j].width = Inches(wd)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.margin_top = cell.margin_bottom = Emu(27432)
            cell.margin_left = cell.margin_right = Emu(73152)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.line_spacing = 1.05
            r = p.add_run(); r.text = val
            if i == 0:
                set_font(r, header_size, True, RGBColor(0xFF, 0xD9, 0xB0))
                cell_fill(cell, '223A63', 92)
            elif i in highlight_rows:
                set_font(r, body_size, True, RED if highlight_color is None else highlight_color)
                cell_fill(cell, 'FF7A59', 16)
            else:
                set_font(r, body_size, False, DARK)
                cell_fill(cell, 'FFFFFF', 5 if i % 2 == 1 else 9)
    return tbl


def picture(slide, path, x, y, max_w, max_h, frame=True):
    """흰 카드 프레임(그림자) 위에 비율 유지 이미지"""
    from PIL import Image
    iw, ih = Image.open(path).size
    scale = min((max_w - 0.24) / iw, (max_h - 0.24) / ih)
    w, h = iw * scale, ih * scale
    px, py = x + (max_w - w) / 2, y + (max_h - h) / 2
    if frame:
        rect(slide, px - 0.12, py - 0.12, w + 0.24, h + 0.24, WHITE)
    return slide.shapes.add_picture(path, Inches(px), Inches(py), Inches(w), Inches(h))
