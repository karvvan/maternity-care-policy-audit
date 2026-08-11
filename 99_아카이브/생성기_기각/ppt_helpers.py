# -*- coding: utf-8 -*-
"""
PPT 공용 헬퍼 v3 — 화이트 테마
순백 배경 위 떠 있는 카드(그림자 원근) + 반투명 그라데이션 블롭 + 3D 조형물 + 삽화.
색 상수 이름은 기존 덱 코드와의 호환을 위해 유지하고 역할만 재해석:
  DARK  = 본문 잉크   NAVY = 카드 제목(네이비)   GRAY = 보조 텍스트
  LIGHT = 흰 카드(마커)   REDBG = 코랄 틴트 경고 카드(마커)   RED = 경고 텍스트
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

import design_assets

ASSETS = design_assets.build_defaults()

# ── 팔레트 (화이트 테마) ────────────────────────────────────
DARK = RGBColor(0x22, 0x2E, 0x40)    # 본문 잉크
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6A, 0x7B, 0x92)    # 보조 텍스트
NAVY = RGBColor(0x17, 0x45, 0x6B)    # 카드 제목·강조 네이비
RED = RGBColor(0xE8, 0x5D, 0x42)     # 경고 코랄
TEAL = RGBColor(0x14, 0x9E, 0x8C)
BLUE = RGBColor(0x3E, 0x6D, 0xB5)
AMBER = RGBColor(0xE8, 0x8A, 0x2E)
LIGHT = RGBColor(0xF2, 0xF6, 0xFA)   # 마커: 흰 카드
REDBG = RGBColor(0xFB, 0xF0, 0xEE)   # 마커: 경고 카드
LINE = RGBColor(0xE2, 0xE9, 0xF1)
INK = DARK
FONT = '맑은 고딕'

_ACCENT1, _ACCENT2 = 'FF7A59', 'FFB86B'
_NAVY1, _NAVY2 = '1B3A5F', '17456B'


def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs, prs.slide_layouts[6]


# ── 저수준 스타일 유틸 ──────────────────────────────────────
def _spPr(sp):
    return sp._element.spPr


def alpha_fill(sp, hex_str, alpha):
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor.from_string(hex_str)
    clr = _spPr(sp).find(qn('a:solidFill')).find(qn('a:srgbClr'))
    a = etree.SubElement(clr, qn('a:alpha'))
    a.set('val', str(int(alpha * 1000)))


def grad_fill(sp, c1, c2, angle=90):
    sp.fill.gradient()
    stops = sp.fill.gradient_stops
    stops[0].color.rgb = RGBColor.from_string(c1)
    stops[0].position = 0.0
    stops[1].color.rgb = RGBColor.from_string(c2)
    stops[1].position = 1.0
    try:
        sp.fill.gradient_angle = angle
    except Exception:
        pass


def line_style(sp, hex_str, alpha=100, width=1.0):
    sp.line.color.rgb = RGBColor.from_string(hex_str)
    sp.line.width = Pt(width)
    if alpha < 100:
        ln = _spPr(sp).find(qn('a:ln'))
        clr = ln.find(qn('a:solidFill')).find(qn('a:srgbClr'))
        a = etree.SubElement(clr, qn('a:alpha'))
        a.set('val', str(int(alpha * 1000)))


def shadow(sp, blur=0.20, dist=0.07, alpha=30, direction=5400000, color='2A3B55'):
    spPr = _spPr(sp)
    old = spPr.find(qn('a:effectLst'))
    if old is not None:
        spPr.remove(old)
    el = etree.SubElement(spPr, qn('a:effectLst'))
    sh = etree.SubElement(el, qn('a:outerShdw'))
    sh.set('blurRad', str(int(Inches(blur))))
    sh.set('dist', str(int(Inches(dist))))
    sh.set('dir', str(direction))
    sh.set('rotWithShape', '0')
    clr = etree.SubElement(sh, qn('a:srgbClr'))
    clr.set('val', color)
    a = etree.SubElement(clr, qn('a:alpha'))
    a.set('val', str(int(alpha * 1000)))


def cell_fill(cell, hex_str, alpha=100):
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ('a:solidFill', 'a:noFill', 'a:gradFill', 'a:pattFill'):
        for e in tcPr.findall(qn(tag)):
            tcPr.remove(e)
    fill = etree.SubElement(tcPr, qn('a:solidFill'))
    clr = etree.SubElement(fill, qn('a:srgbClr'))
    clr.set('val', hex_str)
    if alpha < 100:
        a = etree.SubElement(clr, qn('a:alpha'))
        a.set('val', str(int(alpha * 1000)))


# ── 배경·조형물 ─────────────────────────────────────────────
def bg(slide, kind='main'):
    """순백 배경 + 은은한 블롭 장식 (가장 먼저 호출)"""
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    r.fill.solid(); r.fill.fore_color.rgb = WHITE
    r.line.fill.background(); r.shadow.inherit = False
    if kind == 'title':
        slide.shapes.add_picture(ASSETS['blob_coral'], Inches(7.6), Inches(-3.2), Inches(8.0), Inches(8.0))
        slide.shapes.add_picture(ASSETS['blob_blue'], Inches(-3.4), Inches(2.6), Inches(8.6), Inches(8.6))
        slide.shapes.add_picture(ASSETS['blob_teal'], Inches(4.6), Inches(5.6), Inches(4.6), Inches(4.6))
    elif kind == 'closing':
        slide.shapes.add_picture(ASSETS['blob_blue'], Inches(6.8), Inches(-3.6), Inches(9.0), Inches(9.0))
        slide.shapes.add_picture(ASSETS['blob_coral'], Inches(-3.0), Inches(3.4), Inches(7.6), Inches(7.6))
    else:
        slide.shapes.add_picture(ASSETS['blob_blue'], Inches(-3.8), Inches(4.2), Inches(7.2), Inches(7.2))
        slide.shapes.add_picture(ASSETS['blob_coral'], Inches(9.8), Inches(-3.6), Inches(6.8), Inches(6.8))


def glow(slide, x, y, w, h=0.055):
    slide.shapes.add_picture(ASSETS['glow_bar'], Inches(x), Inches(y), Inches(w), Inches(h))


def deco(slide, key, x, y, w):
    from PIL import Image
    iw, ih = Image.open(ASSETS[key]).size
    slide.shapes.add_picture(ASSETS[key], Inches(x), Inches(y), Inches(w), Inches(w * ih / iw))


# ── 텍스트 ──────────────────────────────────────────────────
def set_font(run, size, bold=False, color=DARK):
    f = run.font
    f.size, f.bold, f.name = Pt(size), bold, FONT
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def text_block(slide, x, y, w, h, items, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (t, size, bold, color, after) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(after)
        p.line_spacing = 1.14
        r = p.add_run(); r.text = t
        set_font(r, size, bold, color)
        if t.startswith('• '):
            pPr = p._p.get_or_add_pPr()
            pPr.set('marL', '160020'); pPr.set('indent', '-160020')
    return tb


# ── 도형 ────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill, line_color=None):
    """테마 인지: LIGHT→떠 있는 흰 카드, REDBG→코랄 틴트 카드, WHITE→흰 액자, 그 외→네이비 그라데이션 블록"""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.adjustments[0] = min(0.5, 0.09 / min(w, h) if min(w, h) > 0.4 else 0.18)
    sp.shadow.inherit = False
    if fill == LIGHT:
        sp.fill.solid(); sp.fill.fore_color.rgb = WHITE
        line_style(sp, 'E4EAF2', 100, 1.0)
        shadow(sp)
    elif fill == REDBG:
        sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xEE)
        line_style(sp, 'F6CDBE', 100, 1.0)
        shadow(sp, alpha=26, color='B24A32')
    elif fill == WHITE:
        sp.fill.solid(); sp.fill.fore_color.rgb = WHITE
        sp.line.fill.background()
        shadow(sp, blur=0.24, alpha=34)
    else:  # 스탯 블록 — 네이비 그라데이션
        grad_fill(sp, _NAVY1, _NAVY2, angle=115)
        sp.line.fill.background()
        shadow(sp, alpha=38)
    return sp


def accent_chip(slide, x, y, w, h):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.adjustments[0] = 0.5
    sp.shadow.inherit = False
    grad_fill(sp, _ACCENT1, _ACCENT2, angle=45)
    sp.line.fill.background()
    shadow(sp, blur=0.12, dist=0.04, alpha=35, color='C24E30')
    return sp


def card(slide, x, y, w, h, title, body, fill=LIGHT, tcolor=NAVY, tsize=13, bsize=11.5):
    rect(slide, x, y, w, h, fill)
    items = [(title, tsize, True, tcolor, 3.5)]
    for line in body:
        items.append((line if line.startswith(('→', '—')) else '• ' + line, bsize, False, DARK, 2))
    text_block(slide, x + 0.2, y + 0.14, w - 0.4, h - 0.28, items)


def header(slide, no, title, footer_label, sub=None):
    """화이트 헤더: 액센트 칩 + 잉크 제목 + 글로우 바 + 페이지 번호"""
    bg(slide, 'main')
    accent_chip(slide, 0.55, 0.52, 0.34, 0.34)
    text_block(slide, 1.05, 0.3, 10.6, 0.75, [(title, 21.5, True, INK, 0)], anchor=MSO_ANCHOR.MIDDLE)
    glow(slide, 0.57, 1.06, 4.6)
    text_block(slide, 11.85, 0.18, 1.3, 0.85, [(f'{no:02d}', 30, True, RGBColor(0xD3, 0xDD, 0xEA), 0)])
    if sub:
        text_block(slide, 0.6, 1.2, 12.2, 0.42, [(sub, 12.5, False, GRAY, 0)])
    text_block(slide, 0.55, 7.1, 12.3, 0.35,
               [(f'2026 명지대학교 창의적 SW프로그램 경진대회 · 빅데이터 분석 — {footer_label}', 8.5, False,
                 RGBColor(0xAB, 0xB7, 0xC7), 0)])


def table(slide, x, y, w, rows, widths=None, header_size=11.5, body_size=11, row_h=0.42,
          highlight_rows=(), highlight_color=None):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(row_h * len(rows)))
    tbl = shape.table
    tbl.first_row = False
    tbl.horz_banding = False
    if widths:
        for j, wd in enumerate(widths):
            tbl.columns[j].width = Inches(wd)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.margin_top = cell.margin_bottom = Emu(27432)
            cell.margin_left = cell.margin_right = Emu(73152)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.line_spacing = 1.05
            r = p.add_run(); r.text = val
            if i == 0:
                set_font(r, header_size, True, WHITE)
                cell_fill(cell, '33393F')   # 무채색 차콜 헤더
            elif i in highlight_rows:
                set_font(r, body_size, True, RED if highlight_color is None else highlight_color)
                cell_fill(cell, 'FFEDE6')
            else:
                set_font(r, body_size, False, DARK)
                cell_fill(cell, 'FFFFFF' if i % 2 == 1 else 'F5F8FC')
    return tbl


def picture(slide, path, x, y, max_w, max_h, frame=True):
    from PIL import Image
    iw, ih = Image.open(path).size
    scale = min((max_w - 0.24) / iw, (max_h - 0.24) / ih)
    w, h = iw * scale, ih * scale
    px, py = x + (max_w - w) / 2, y + (max_h - h) / 2
    if frame:
        rect(slide, px - 0.12, py - 0.12, w + 0.24, h + 0.24, WHITE)
    return slide.shapes.add_picture(path, Inches(px), Inches(py), Inches(w), Inches(h))
